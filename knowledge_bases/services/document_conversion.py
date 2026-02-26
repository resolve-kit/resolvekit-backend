from __future__ import annotations

from dataclasses import dataclass
import io
from pathlib import Path
import os
from tempfile import NamedTemporaryFile
from typing import Iterable

import chardet

from knowledge_bases.config import settings

try:  # pragma: no cover - optional runtime dependency
    from markitdown import MarkItDown
except Exception:  # pragma: no cover - optional runtime dependency
    MarkItDown = None  # type: ignore[assignment]


_DEFAULT_ALLOWED_EXTENSIONS = (
    ".txt",
    ".md",
    ".markdown",
    ".pdf",
    ".doc",
    ".docx",
    ".ppt",
    ".pptx",
    ".rtf",
    ".odt",
    ".html",
    ".htm",
    ".csv",
    ".tsv",
    ".xlsx",
    ".xls",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
)

_TEXT_FALLBACK_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".html",
    ".htm",
    ".csv",
    ".tsv",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
}


@dataclass(slots=True)
class ConvertedUpload:
    filename: str
    extension: str
    title: str
    markdown: str


class DocumentConversionError(ValueError):
    pass


def parse_allowed_extensions(raw: str | None) -> set[str]:
    if not raw:
        return set(_DEFAULT_ALLOWED_EXTENSIONS)

    extensions: set[str] = set()
    for token in raw.split(","):
        ext = token.strip().lower()
        if not ext:
            continue
        if not ext.startswith("."):
            ext = f".{ext}"
        extensions.add(ext)

    return extensions or set(_DEFAULT_ALLOWED_EXTENSIONS)


def _normalize_filename(filename: str) -> str:
    cleaned = (filename or "").strip()
    if cleaned:
        return cleaned
    raise DocumentConversionError("A filename is required for upload")


def _extension_for(filename: str) -> str:
    return Path(filename).suffix.lower()


def _select_title(title: str | None, filename: str) -> str:
    cleaned = (title or "").strip()
    if cleaned:
        return cleaned[:255]

    stem = Path(filename).stem.strip()
    if stem:
        return stem[:255]

    return "Uploaded File"


def _decode_text(content: bytes) -> str:
    if not content:
        return ""

    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        detected = chardet.detect(content)
        encoding = (detected.get("encoding") or "utf-8").strip() or "utf-8"
        try:
            return content.decode(encoding, errors="replace")
        except Exception:
            return content.decode("utf-8", errors="replace")


def _extract_text_from_markitdown_result(result: object) -> str:
    if isinstance(result, str):
        return result.strip()

    for attr in ("text_content", "markdown", "text", "content"):
        value = getattr(result, attr, None)
        if isinstance(value, str) and value.strip():
            return value.strip()

    return ""


def _extract_with_markitdown(filename: str, content: bytes) -> str:
    if MarkItDown is None:
        return ""

    temp_path: str | None = None
    try:
        suffix = _extension_for(filename) or ".txt"
        with NamedTemporaryFile("wb", suffix=suffix, delete=False) as temp:
            temp.write(content)
            temp_path = temp.name

        converter = MarkItDown()
        result = converter.convert(temp_path)
        return _extract_text_from_markitdown_result(result)
    except Exception:
        return ""
    finally:
        if temp_path:
            try:
                os.unlink(temp_path)
            except OSError:
                pass


def _extract_pdf(content: bytes) -> str:
    try:  # pragma: no cover - exercised only when dependency exists
        from pypdf import PdfReader
    except Exception:
        return ""

    try:
        reader = PdfReader(io.BytesIO(content))
        chunks = [page.extract_text() or "" for page in reader.pages]
    except Exception:
        return ""

    return "\n\n".join(piece.strip() for piece in chunks if piece and piece.strip()).strip()


def _extract_docx(content: bytes) -> str:
    try:  # pragma: no cover - exercised only when dependency exists
        from docx import Document
    except Exception:
        return ""

    try:
        document = Document(io.BytesIO(content))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    except Exception:
        return ""

    return "\n\n".join(paragraphs).strip()


def _extract_pptx(content: bytes) -> str:
    try:  # pragma: no cover - exercised only when dependency exists
        from pptx import Presentation
    except Exception:
        return ""

    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception:
        return ""

    text_parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            value = getattr(shape, "text", None)
            if isinstance(value, str) and value.strip():
                text_parts.append(value.strip())

    return "\n\n".join(text_parts).strip()


def _extract_rtf(content: bytes) -> str:
    try:  # pragma: no cover - exercised only when dependency exists
        from striprtf.striprtf import rtf_to_text
    except Exception:
        return ""

    raw = _decode_text(content)
    if not raw.strip():
        return ""

    try:
        return rtf_to_text(raw).strip()
    except Exception:
        return ""


def _extract_with_fallback(extension: str, content: bytes) -> str:
    if extension in _TEXT_FALLBACK_EXTENSIONS:
        return _decode_text(content).strip()
    if extension == ".pdf":
        return _extract_pdf(content)
    if extension == ".docx":
        return _extract_docx(content)
    if extension == ".pptx":
        return _extract_pptx(content)
    if extension == ".rtf":
        return _extract_rtf(content)
    return ""


def _extract_with_ocr(content: bytes) -> str:
    try:  # pragma: no cover - OCR support is optional
        import pytesseract  # noqa: F401
    except Exception as exc:
        raise DocumentConversionError("OCR is enabled but OCR dependencies are not installed") from exc

    # OCR backend is feature-flagged and intentionally opt-in.
    # If dependencies are installed, this can be replaced with full OCR extraction.
    return ""


def convert_uploaded_file_bytes(
    *,
    filename: str,
    content: bytes,
    content_type: str | None,
    title: str | None,
    max_file_bytes: int | None = None,
    allowed_extensions: Iterable[str] | None = None,
    ocr_enabled: bool | None = None,
) -> ConvertedUpload:
    normalized_filename = _normalize_filename(filename)
    extension = _extension_for(normalized_filename)
    normalized_extensions = (
        {ext.lower() for ext in allowed_extensions}
        if allowed_extensions is not None
        else parse_allowed_extensions(settings.upload_allowed_extensions)
    )

    if not extension or extension not in normalized_extensions:
        allowed = ", ".join(sorted(normalized_extensions))
        raise DocumentConversionError(
            f"Unsupported file extension '{extension or '<none>'}'. Allowed: {allowed}"
        )

    file_limit = max_file_bytes if max_file_bytes is not None else settings.upload_max_file_bytes
    file_size = len(content)
    if file_size == 0:
        raise DocumentConversionError("Uploaded file is empty")
    if file_limit > 0 and file_size > file_limit:
        raise DocumentConversionError(
            f"Uploaded file size exceeds limit ({file_size} bytes > {file_limit} bytes)"
        )

    extracted = _extract_with_markitdown(normalized_filename, content)
    if not extracted:
        extracted = _extract_with_fallback(extension, content)

    should_try_ocr = settings.upload_ocr_enabled if ocr_enabled is None else ocr_enabled
    if not extracted and should_try_ocr:
        extracted = _extract_with_ocr(content)

    normalized_markdown = extracted.strip()
    if not normalized_markdown:
        raise DocumentConversionError(
            "No extractable text found in uploaded file. "
            "If this is a scanned document, enable OCR support in KB service settings."
        )

    return ConvertedUpload(
        filename=normalized_filename,
        extension=extension,
        title=_select_title(title, normalized_filename),
        markdown=normalized_markdown,
    )
